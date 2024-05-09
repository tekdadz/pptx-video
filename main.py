import aspose.slides as slides
import aspose.pydrawing as draw
from pptx import Presentation
import eyed3
import subprocess

origin_path = "demo.pptx"
export_path = "export.pptx"
video_export_path = "export.mp4"
audio_files = ["slide1.mp3", "slide2.mp3", "slide3.mp3", "slide4.mp3", "slide5.mp3"]
durations = []
slide_images = []
image_width = 1920
image_height = 1080

with slides.Presentation(origin_path) as presentation:

    for index, audio_file_path in enumerate(audio_files):
        sld = presentation.slides[index]
        duration = eyed3.load(audio_file_path).info.time_secs
        durations.append(duration)

        print(audio_file_path, index)
        print (duration)

        with open(audio_file_path, "rb") as in_file:
            audio_frame = sld.shapes.add_audio_frame_embedded(50, 150, 100, 100, in_file)
            audio_frame.play_mode = slides.AudioPlayModePreset.AUTO
            audio_frame.hide_at_showing = True
            audio_frame.volume = slides.AudioVolumeMode.LOUD
        sld.slide_show_transition.advance_after_time = int(duration) * 1000

    presentation.save(export_path, slides.export.SaveFormat.PPTX)

pre = Presentation(export_path)
for slide in pre.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            if "Aspose Pty Ltd." in shape.text:
                shape.text = ""
pre.save(export_path)


# def export_slides_as_images(ppt_path, output_folder):
#     pres = slides.Presentation(ppt_path)
#     for i, slide in enumerate(pres.slides):
#         slide_image_path = f"{output_folder}/slide_{i + 1}.png"
#         slide_images.append(slide_image_path)
#         slide.get_image(draw.Size(image_width, image_height)).save(slide_image_path, slides.ImageFormat.PNG)


# output_folder = 'images'
# export_slides_as_images(origin_path, output_folder)


# def create_video(slide_images, audio_files, durations, output_video):
#     input_args = []
#     filter_complex = []
#     filter_complex_sub = []
#     filter_complex_sub2 = []
#     duration_total = 0

#     # ffmpeg -loop 1 -t 44 -i slide1.png -i audio1.mp3 -loop 1 -t 21 -i slide2.png -i audio2.mp3 -loop 1 -t 31 -i slide3.png -i audio3.mp3 
#     # -filter_complex "[0:v]scale=1920:1080,setpts=PTS-STARTPTS,trim=duration=44[v0]; [2:v]scale=1920:1080,setpts=PTS-STARTPTS,trim=duration=21[v1]; [4:v]scale=1920:1080,setpts=PTS-STARTPTS,trim=duration=31[v2]; 
#     # [v0][1:a][v1][3:a][v2][5:a]concat=n=3:v=1:a=1[v][a]" -map "[v]" -map "[a]" -c:v libx264 -tune stillimage -c:a aac -b:a 192k -pix_fmt yuv420p output.mp4
#     for i, (audio, duration, slide_image_path) in enumerate(zip(audio_files, durations, slide_images)):
#         input_args.extend(['-loop', '1', '-t', str(duration), '-i', slide_image_path])  # Image input
#         input_args.extend(['-i', audio])  # Audio input
#         filter_complex.append(f"[{2*i}:v]scale={image_width}:{image_height},setpts=PTS-STARTPTS,trim=duration={duration}[v{i}]")
#         filter_complex_sub.append(f"[v{i}][{2*i+1}:a]")
#         duration_total += duration

#     filter_complex_sub2 = f"concat=n={len(audio_files)}:v=1:a=1[v][a]"

#     output_args = ['-map', '"[v]"', '-map', '"[a]"', '-c:v libx264 -tune stillimage -c:a aac -b:a 192k -pix_fmt yuv420p', output_video]

#     # Construct and run the ffmpeg command
#     command = ['ffmpeg'] + input_args + ['-filter_complex', '"' + ';'.join(filter_complex) + ';' + ''.join(filter_complex_sub) + filter_complex_sub2 + '"'] + output_args

#     print(command)

#     try:
#         subprocess.run(command, check=True)
#         print("Video created successfully")
#     except subprocess.CalledProcessError as e:
#         print(f"Error: {e}")


# output_video = 'output_video.mp4'

# create_video(slide_images, audio_files, durations, output_video)